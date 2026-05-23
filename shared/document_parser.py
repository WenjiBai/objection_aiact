"""Upload parsers — turns Streamlit UploadedFile into clean UTF-8 text.

Fixes the bug where binary PDF/DOCX/PPTX bytes were decoded as UTF-8, producing
hundreds of thousands of replacement characters and blowing the Anthropic
30k-token/min rate limit on the first agent call.

Dispatches by mime_type (preferred) with filename-extension fallback.
Truncates oversized documents to protect the LLM token budget.
"""

from __future__ import annotations

import io
from dataclasses import dataclass


# Cap per-document content. ~30k characters ≈ ~7.5k tokens; with 6 agents and
# RAG context this keeps a single hearing well under per-minute limits.
MAX_CHARS_PER_DOC = 30_000


@dataclass
class ParsedDocument:
    text: str
    truncated: bool
    original_chars: int
    parser: str  # "pdf" | "docx" | "pptx" | "text" | "fallback"


# ---------------------------------------------------------------- per-format

def _parse_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages: list[str] = []
    for page in reader.pages:
        try:
            pages.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n\n".join(p.strip() for p in pages if p and p.strip())


def _parse_docx(data: bytes) -> str:
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(data))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation

    pres = Presentation(io.BytesIO(data))
    out: list[str] = []
    for i, slide in enumerate(pres.slides, start=1):
        slide_lines = [f"# Slide {i}"]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    slide_lines.append(text)
        if len(slide_lines) > 1:
            out.append("\n".join(slide_lines))
    return "\n\n".join(out)


def _parse_text(data: bytes) -> str:
    # Real text files: utf-8 with replacement is fine; binary garbage never
    # reaches here because we dispatch by mime/extension first.
    return data.decode("utf-8", errors="replace")


# ---------------------------------------------------------------- dispatch

def _kind(filename: str, mime: str | None) -> str:
    mime = (mime or "").lower()
    name = filename.lower()
    if "pdf" in mime or name.endswith(".pdf"):
        return "pdf"
    if "wordprocessingml" in mime or name.endswith(".docx"):
        return "docx"
    if "presentationml" in mime or name.endswith(".pptx"):
        return "pptx"
    if name.endswith((".txt", ".md", ".markdown")) or mime.startswith("text/"):
        return "text"
    return "unknown"


def extract_text(uploaded_file) -> ParsedDocument:
    """Extract clean UTF-8 text from a Streamlit UploadedFile.

    Never raises on parser failure — falls back to a placeholder string so the
    courtroom pipeline keeps running.
    """
    data: bytes = uploaded_file.getvalue()
    kind = _kind(uploaded_file.name, getattr(uploaded_file, "type", None))

    parser_used = kind
    try:
        if kind == "pdf":
            text = _parse_pdf(data)
        elif kind == "docx":
            text = _parse_docx(data)
        elif kind == "pptx":
            text = _parse_pptx(data)
        elif kind == "text":
            text = _parse_text(data)
        else:
            # Unknown format — try utf-8 but don't poison context if it's binary.
            decoded = data.decode("utf-8", errors="ignore")
            text = decoded if decoded.strip() else f"[Unsupported file format: {uploaded_file.name}]"
            parser_used = "fallback"
    except Exception as exc:
        text = f"[Failed to parse {uploaded_file.name}: {type(exc).__name__}: {exc}]"
        parser_used = "fallback"

    original = len(text)
    truncated = original > MAX_CHARS_PER_DOC
    if truncated:
        text = text[:MAX_CHARS_PER_DOC] + f"\n\n[... truncated — {original - MAX_CHARS_PER_DOC} more characters omitted to protect token budget ...]"

    return ParsedDocument(
        text=text,
        truncated=truncated,
        original_chars=original,
        parser=parser_used,
    )
